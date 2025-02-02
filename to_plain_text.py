import sys
import pypandoc
import PyPDF2
from pptx import Presentation

def todo():
    return print_return_error("Not yet implemented.")

def print_return_error(err):
    print(err)
    return err

def with_pandoc(filepath):
    try:
        plain_text = pypandoc.convert_file(filepath, 'plain')
        return plain_text
        
    except Exception as e:
        return print_return_error(f"Error converting {filepath} to plain text.")


def from_docx(filepath):
    return with_pandoc(filepath)

def from_doc(filepath):
    return with_pandoc(filepath)

def from_pptx(filepath):
    presentation = Presentation(filepath)
    extracted_text = []
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text + "\n"
        extracted_text.append(slide_text)
    return "\n".join(extracted_text)


def from_pdf(filepath):
    try:
        with open(filepath, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            num_pages = len(reader.pages) 
            
            plain_text = ''
            for page_num in range(num_pages):
                page = reader.pages[page_num]
                plain_text += page.extract_text()
                
        return plain_text.strip()
        
    except Exception as e:
        return print_return_error(f"Error converting {filepath} to plain text.")

def from_plain_text(filepath):
    with open(filepath, "r") as f:
        output = f.read()
    
    if output is None or output == "":
        return print_return_error(f"Error gathering text from {filepath}.")

    return output

def get_content(filepath):
    print(filepath)
    parts = filepath.split(".")
    ext = parts[len(parts) - 1]

    if ext == "txt":
        return from_plain_text(filepath)
    elif ext == "pdf":
        return from_pdf(filepath)
    elif ext == "docx":
        return from_docx(filepath)
    elif ext == "doc":
        return from_doc(filepath)
    elif ext == "pptx":
        return from_pptx(filepath)
    else:
        return print_return_error(f"Extension {ext} is not yet supported.")
